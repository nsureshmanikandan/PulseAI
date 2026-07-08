"""
PulseAI - Enterprise Presentation Generator
Professional dark-theme PPT for Customer & Leadership briefings
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
from pptx.enum.dml import MSO_THEME_COLOR
import copy

# ── Color Palette ──────────────────────────────────────────────────────────────
NAVY        = RGBColor(0x0A, 0x0F, 0x1E)   # #0a0f1e  background
DARK_CARD   = RGBColor(0x0F, 0x17, 0x2A)   # #0f172a  card bg
BLUE        = RGBColor(0x3B, 0x82, 0xF6)   # #3b82f6  primary
INDIGO      = RGBColor(0x63, 0x66, 0xF1)   # #6366f1  accent
EMERALD     = RGBColor(0x10, 0xB9, 0x81)   # #10b981  success
AMBER       = RGBColor(0xF5, 0x9E, 0x0B)   # #f59e0b  warning
RED         = RGBColor(0xEF, 0x44, 0x44)   # #ef4444  danger
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
GRAY_300    = RGBColor(0xCB, 0xD5, 0xE1)   # #cbd5e1
GRAY_500    = RGBColor(0x64, 0x74, 0x8B)   # #64748b
SLATE_800   = RGBColor(0x1E, 0x29, 0x3B)   # #1e293b

prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)

BLANK = prs.slide_layouts[6]  # completely blank

# ── Helpers ────────────────────────────────────────────────────────────────────

def slide():
    s = prs.slides.add_slide(BLANK)
    bg = s.background.fill
    bg.solid()
    bg.fore_color.rgb = NAVY
    return s

def rect(s, x, y, w, h, fill=DARK_CARD, alpha=None):
    shape = s.shapes.add_shape(1, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.fill.background()
    return shape

def txt(s, text, x, y, w, h, size=18, bold=False, color=WHITE,
        align=PP_ALIGN.LEFT, italic=False, wrap=True):
    tb = s.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tb.word_wrap = wrap
    tf = tb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return tb

def accent_bar(s, x, y, w=0.05, h=0.45, color=BLUE):
    rect(s, x, y, w, h, fill=color)

def badge(s, text, x, y, w, h, bg=BLUE, fg=WHITE, size=10):
    r = rect(s, x, y, w, h, fill=bg)
    txt(s, text, x+0.05, y+0.02, w-0.1, h-0.04, size=size, bold=True,
        color=fg, align=PP_ALIGN.CENTER)
    return r

def divider(s, y, color=BLUE):
    line = s.shapes.add_shape(1, Inches(0.5), Inches(y), Inches(12.33), Inches(0.02))
    line.fill.solid()
    line.fill.fore_color.rgb = color
    line.line.fill.background()

def bullet_slide_body(s, items, x, y, col_color=BLUE):
    for i, (icon, title, body) in enumerate(items):
        iy = y + i * 0.72
        rect(s, x, iy, 11.8, 0.60, fill=SLATE_800)
        rect(s, x, iy, 0.04, 0.60, fill=col_color)
        txt(s, icon, x+0.08, iy+0.08, 0.4, 0.45, size=16)
        txt(s, title, x+0.52, iy+0.04, 3.2, 0.28, size=11, bold=True, color=WHITE)
        txt(s, body,  x+0.52, iy+0.30, 11.0, 0.28, size=9.5, color=GRAY_300)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 1 — Title / Hero
# ═══════════════════════════════════════════════════════════════════════════════
s1 = slide()

# Gradient overlay left panel
rect(s1, 0, 0, 6.5, 7.5, fill=RGBColor(0x06, 0x0B, 0x18))
rect(s1, 0, 0, 0.08, 7.5, fill=BLUE)

# Logo badge
rect(s1, 0.55, 0.6, 0.9, 0.9, fill=BLUE)
txt(s1, "P", 0.55, 0.6, 0.9, 0.9, size=32, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

txt(s1, "PulseAI", 1.6, 0.65, 4.5, 0.7, size=32, bold=True, color=WHITE)
txt(s1, "Data Analytics", 1.6, 1.28, 4.5, 0.4, size=13, color=BLUE)

# Tagline
txt(s1, "Turn Raw Data Into\nDecisions — Instantly.",
    0.55, 2.1, 5.7, 1.4, size=26, bold=True, color=WHITE)

txt(s1, "Enterprise-grade AI analytics platform that transforms\nExcel & CSV files into executive dashboards, AI-powered\ninsights, and interactive data stories — in minutes.",
    0.55, 3.65, 5.7, 1.5, size=11.5, color=GRAY_300)

# Stat pills bottom-left
for i, (val, lbl) in enumerate([("3,411+", "Rows Analyzed"), ("5", "Excel Tabs"), ("GPT-4o", "AI Engine")]):
    bx = 0.55 + i * 1.9
    rect(s1, bx, 5.5, 1.7, 0.7, fill=SLATE_800)
    rect(s1, bx, 5.5, 0.04, 0.7, fill=BLUE)
    txt(s1, val, bx+0.12, 5.52, 1.5, 0.32, size=14, bold=True, color=BLUE)
    txt(s1, lbl, bx+0.12, 5.82, 1.5, 0.28, size=8.5, color=GRAY_300)

# Right panel — mock dashboard preview blocks
for ry, rh, rc in [(0.3,0.9,SLATE_800),(1.35,0.5,SLATE_800),(1.35,0.5,SLATE_800),
                    (2.0,2.2,SLATE_800),(4.4,2.7,SLATE_800)]:
    pass  # replaced by visual cards below

# Right side visual cards
rect(s1, 6.8, 0.25, 6.1, 7.0, fill=RGBColor(0x06, 0x0D, 0x20))
rect(s1, 6.8, 0.25, 0.04, 7.0, fill=INDIGO)

txt(s1, "LIVE DEMO SNAPSHOT", 7.1, 0.35, 5.5, 0.35, size=8, color=INDIGO, bold=True)

# KPI row
for i, (v, l, c) in enumerate([("$743.9M","Total Outstanding",BLUE),("34%","Default+Watch Rate",RED),("657","Avg Credit Score",EMERALD)]):
    cx = 7.1 + i*1.95
    rect(s1, cx, 0.8, 1.8, 0.9, fill=SLATE_800)
    rect(s1, cx, 0.8, 0.04, 0.9, fill=c)
    txt(s1, v, cx+0.12, 0.82, 1.6, 0.42, size=17, bold=True, color=c)
    txt(s1, l, cx+0.12, 1.22, 1.6, 0.35, size=8.5, color=GRAY_300)

# Chart preview blocks
rect(s1, 7.1, 1.9, 5.6, 2.0, fill=SLATE_800)
txt(s1, "Distribution of Principal Amount", 7.2, 1.95, 5.4, 0.28, size=9, bold=True, color=WHITE)
for i in range(12):
    bh = 0.2 + (i % 5) * 0.22
    rect(s1, 7.2 + i*0.44, 3.65 - bh, 0.36, bh,
         fill=BLUE if i % 3 == 0 else INDIGO if i % 3 == 1 else RGBColor(0x06,0xB6,0xD4))

rect(s1, 7.1, 4.1, 5.6, 2.2, fill=SLATE_800)
txt(s1, "Loan Status Breakdown", 7.2, 4.15, 5.4, 0.28, size=9, bold=True, color=WHITE)
for i, (pct, c, lbl) in enumerate([(60,BLUE,"PAIDOFF"),(20,EMERALD,"COLLECTION"),(20,AMBER,"COL_PAIDOFF")]):
    rect(s1, 7.2, 4.55+i*0.45, pct/100*4.5, 0.32, fill=c)
    txt(s1, f"{lbl}  {pct}%", 7.2, 4.55+i*0.45, 5.0, 0.32, size=8.5, color=WHITE)

txt(s1, "Powered by GPT-4o  |  Accenture GenAI Gateway  |  Built on React + FastAPI",
    0.55, 7.1, 12.0, 0.3, size=8, color=GRAY_500, align=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 2 — Problem Statement
# ═══════════════════════════════════════════════════════════════════════════════
s2 = slide()
rect(s2, 0, 0, 13.33, 1.1, fill=RGBColor(0x06,0x0B,0x18))
rect(s2, 0, 0, 0.06, 1.1, fill=RED)
txt(s2, "The Problem", 0.35, 0.12, 8, 0.5, size=24, bold=True, color=WHITE)
txt(s2, "Why organizations struggle with their data today", 0.35, 0.62, 9, 0.35, size=12, color=GRAY_300)
badge(s2, "CHALLENGE", 10.8, 0.25, 1.8, 0.45, bg=RED)

problems = [
    ("📊", "Data Trapped in Spreadsheets",
     "Finance teams spend 70% of their time manually formatting Excel files rather than making decisions from the data."),
    ("🤖", "AI Benefits Locked Behind Technical Barriers",
     "Data science queries, Python notebooks, and SQL skills are required just to answer a simple business question."),
    ("⏱️", "Slow Time-to-Insight",
     "Executive reports take days to prepare. By the time insights reach leadership, the window for action has passed."),
    ("🔗", "No Cross-Tab Relationship Visibility",
     "Multi-sheet Excel workbooks contain hidden relationships between data — but no tool surfaces them automatically."),
    ("📉", "Generic Dashboards, Not Business Context",
     "Existing BI tools show charts but cannot explain what the data means in plain English for non-technical stakeholders."),
]
bullet_slide_body(s2, problems, 0.5, 1.25, col_color=RED)

rect(s2, 0.5, 6.7, 12.3, 0.55, fill=RGBColor(0x3B,0x0A,0x0A))
txt(s2, "Result: Missed opportunities, delayed decisions, and a growing gap between data availability and business action.",
    0.65, 6.75, 12.0, 0.42, size=10.5, bold=True, color=RGBColor(0xFC,0xA5,0xA5))


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 3 — Solution Overview
# ═══════════════════════════════════════════════════════════════════════════════
s3 = slide()
rect(s3, 0, 0, 13.33, 1.1, fill=RGBColor(0x06,0x0B,0x18))
rect(s3, 0, 0, 0.06, 1.1, fill=EMERALD)
txt(s3, "Introducing PulseAI", 0.35, 0.12, 9, 0.5, size=24, bold=True, color=WHITE)
txt(s3, "One platform. Any data file. Instant intelligence.", 0.35, 0.62, 9, 0.35, size=12, color=GRAY_300)
badge(s3, "SOLUTION", 10.8, 0.25, 1.8, 0.45, bg=EMERALD)

# Flow diagram
steps = [
    (BLUE,   "01", "UPLOAD",    "Drag & drop Excel\nor CSV files"),
    (INDIGO, "02", "ANALYZE",   "Auto schema detect\n+ relationship map"),
    (EMERALD,"03", "VISUALIZE", "6+ auto-charts\n+ KPI dashboard"),
    (AMBER,  "04", "ASK AI",    "Natural language\nquestions & tables"),
    (RED,    "05", "DECIDE",    "Actionable insights\nin plain English"),
]
for i, (col, num, title, body) in enumerate(steps):
    bx = 0.4 + i * 2.55
    rect(s3, bx, 1.3, 2.3, 3.2, fill=SLATE_800)
    rect(s3, bx, 1.3, 2.3, 0.06, fill=col)
    txt(s3, num, bx+0.1, 1.45, 2.1, 0.5, size=26, bold=True, color=col)
    txt(s3, title, bx+0.1, 1.95, 2.1, 0.4, size=12, bold=True, color=WHITE)
    txt(s3, body, bx+0.1, 2.45, 2.1, 0.9, size=9.5, color=GRAY_300)
    if i < 4:
        txt(s3, "->", bx+2.35, 2.5, 0.2, 0.4, size=16, bold=True, color=col)

# Bottom value prop
for i, (icon, val) in enumerate([
    ("< 2 min", "Time to first insight after upload"),
    ("Zero code", "No SQL, Python, or BI training needed"),
    ("GPT-4o", "Powered by frontier AI for accurate analysis"),
    ("Any data", "Finance, HR, Sales, Operations — all domains"),
]):
    bx = 0.4 + i * 3.2
    rect(s3, bx, 4.9, 3.0, 1.0, fill=RGBColor(0x0F,0x17,0x2A))
    rect(s3, bx, 4.9, 0.04, 1.0, fill=BLUE)
    txt(s3, icon, bx+0.15, 4.95, 2.8, 0.4, size=14, bold=True, color=BLUE)
    txt(s3, val,  bx+0.15, 5.35, 2.8, 0.5, size=9,  color=GRAY_300)

txt(s3, "PulseAI bridges the gap between raw data files and confident business decisions — for every user, not just data scientists.",
    0.5, 6.1, 12.3, 0.45, size=10.5, color=GRAY_300, align=PP_ALIGN.CENTER, italic=True)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 4 — Architecture
# ═══════════════════════════════════════════════════════════════════════════════
s4 = slide()
rect(s4, 0, 0, 13.33, 1.1, fill=RGBColor(0x06,0x0B,0x18))
rect(s4, 0, 0, 0.06, 1.1, fill=INDIGO)
txt(s4, "System Architecture", 0.35, 0.12, 9, 0.5, size=24, bold=True, color=WHITE)
txt(s4, "Modern cloud-native stack — React frontend + FastAPI backend + GPT-4o AI layer", 0.35, 0.62, 10, 0.35, size=11, color=GRAY_300)
badge(s4, "ARCHITECTURE", 10.5, 0.25, 2.1, 0.45, bg=INDIGO)

layers = [
    (BLUE,   "FRONTEND  (React 18 + Vite + TypeScript)",
     ["Executive Dashboard", "Analyst Workbench", "AI Chat (WebSocket)", "Data Sources Manager"],
     ["Zustand state", "Plotly.js charts", "ReactMarkdown", "Tailwind CSS"]),
    (INDIGO, "BACKEND API  (FastAPI + Python 3.13)",
     ["/api/datasets", "/api/analytics", "/api/ai", "/ws/chat"],
     ["SQLAlchemy ORM", "Uvicorn ASGI", "Pydantic validation", "CORS middleware"]),
    (EMERALD,"AI ENGINE  (Accenture GenAI Gateway)",
     ["GPT-4o LLM", "Narrative Generator", "Insight Engine", "Suggested Questions"],
     ["AccentureLbpassLLM", "lbpass API key", "Azure OpenAI endpoint", "Streaming support"]),
    (AMBER,  "DATA LAYER  (SQLite + Local Storage)",
     ["Dataset registry", "Tab profiles", "Relationships", "Chat sessions"],
     ["openpyxl parser", "pandas DataFrames", "Schema detector", "Relationship finder"]),
]

for i, (col, title, left, right) in enumerate(layers):
    iy = 1.25 + i * 1.48
    rect(s4, 0.4, iy, 12.5, 1.3, fill=SLATE_800)
    rect(s4, 0.4, iy, 0.06, 1.3, fill=col)
    txt(s4, title, 0.6, iy+0.08, 5.5, 0.3, size=10, bold=True, color=col)
    for j, item in enumerate(left):
        txt(s4, f"+ {item}", 0.6, iy+0.42+j*0.2, 3.5, 0.22, size=8.5, color=WHITE)
    for j, item in enumerate(right):
        txt(s4, f"| {item}", 4.5, iy+0.42+j*0.2, 3.5, 0.22, size=8.5, color=GRAY_300)

# Right side: data flow arrow
rect(s4, 12.4, 1.25, 0.55, 5.92, fill=RGBColor(0x06,0x0B,0x18))
txt(s4, "DATA FLOW", 12.35, 3.6, 0.6, 0.6, size=7, bold=True, color=GRAY_500)
for i in range(3):
    txt(s4, "v", 12.45, 2.35 + i*1.48, 0.4, 0.3, size=14, bold=True, color=BLUE)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 5 — Feature: Executive Dashboard
# ═══════════════════════════════════════════════════════════════════════════════
s5 = slide()
rect(s5, 0, 0, 13.33, 1.1, fill=RGBColor(0x06,0x0B,0x18))
rect(s5, 0, 0, 0.06, 1.1, fill=BLUE)
txt(s5, "Feature 1 — Executive Dashboard", 0.35, 0.12, 9, 0.5, size=22, bold=True, color=WHITE)
txt(s5, "Auto-generated KPIs, AI narrative, and 6 chart types from any uploaded dataset", 0.35, 0.62, 10, 0.35, size=11, color=GRAY_300)
badge(s5, "LIVE FEATURE", 10.5, 0.25, 2.1, 0.45, bg=BLUE)

# Left: mock dashboard
rect(s5, 0.4, 1.2, 7.5, 5.9, fill=RGBColor(0x06,0x0D,0x20))
txt(s5, "Banking_LoanPortfolio.xlsx  |  Loans tab", 0.55, 1.3, 7.0, 0.3, size=8.5, color=GRAY_500)

# KPI cards
for i, (v, l, c) in enumerate([("$743.9M","Total Outstanding",BLUE),("34%","Default+Watch Rate",RED),("657","Avg Credit Score",EMERALD),("80 days","Avg DPD",AMBER)]):
    cx = 0.55 + i * 1.8
    rect(s5, cx, 1.7, 1.65, 0.85, fill=SLATE_800)
    txt(s5, v, cx+0.1, 1.74, 1.45, 0.36, size=14, bold=True, color=c)
    txt(s5, l, cx+0.1, 2.08, 1.45, 0.32, size=7.5, color=GRAY_300)

# AI narrative card
rect(s5, 0.55, 2.72, 7.2, 0.75, fill=RGBColor(0x0F,0x17,0x2A))
rect(s5, 0.55, 2.72, 0.04, 0.75, fill=INDIGO)
txt(s5, "AI Narrative:", 0.68, 2.76, 2.0, 0.28, size=8, bold=True, color=INDIGO)
txt(s5, '"Banking_LoanPortfolio.xlsx / Loans shows a Default+Watch Rate of 34.0% with $743.9M total outstanding. Average credit score is 657 and average DPD is 80 days. Outliers detected in outstanding_balance and EMI — these may warrant risk escalation."',
    0.68, 3.02, 7.0, 0.38, size=7.5, color=GRAY_300, italic=True)

# Charts row
for i, (t, c) in enumerate([("Principal Distribution", BLUE), ("Loan Status Breakdown", EMERALD)]):
    cx = 0.55 + i * 3.65
    rect(s5, cx, 3.6, 3.4, 1.4, fill=SLATE_800)
    txt(s5, t, cx+0.1, 3.65, 3.2, 0.28, size=8, bold=True, color=WHITE)
    if i == 0:
        for j in range(8):
            bh = 0.15 + (j%4)*0.18
            rect(s5, cx+0.15+j*0.38, 4.77-bh, 0.32, bh, fill=c)
    else:
        for j, (p, cl) in enumerate([(60,BLUE),(20,EMERALD),(20,AMBER)]):
            rect(s5, cx+0.15, 3.98+j*0.32, p/100*3.0, 0.26, fill=cl)

# Right: features list
feats = [
    ("KPI Cards", "Smart KPI engine: rate KPIs from status columns, currency totals, score averages — color-coded red/green/blue"),
    ("AI Narrative", "GPT-4o writes a plain-English paragraph describing the dataset"),
    ("6 Chart Types", "Histogram, scatter, heatmap, box plots, bar, donut — all automatic"),
    ("Outlier Alerts", "Flags anomalies and statistical warnings in the data"),
    ("Tab Switching", "Navigate between Excel tabs — each gets its own full dashboard"),
    ("Animated Loading", "Enterprise-grade shimmer skeletons and progress bar"),
]
txt(s5, "WHAT YOU GET", 8.2, 1.22, 4.8, 0.3, size=8.5, bold=True, color=BLUE)
for i, (t, d) in enumerate(feats):
    rect(s5, 8.2, 1.6+i*0.82, 4.8, 0.72, fill=SLATE_800)
    rect(s5, 8.2, 1.6+i*0.82, 0.04, 0.72, fill=BLUE if i%2==0 else INDIGO)
    txt(s5, t, 8.32, 1.62+i*0.82, 4.6, 0.3, size=9.5, bold=True, color=WHITE)
    txt(s5, d, 8.32, 1.9+i*0.82,  4.6, 0.3, size=8,   color=GRAY_300)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 6 — Feature: AI Chat
# ═══════════════════════════════════════════════════════════════════════════════
s6 = slide()
rect(s6, 0, 0, 13.33, 1.1, fill=RGBColor(0x06,0x0B,0x18))
rect(s6, 0, 0, 0.06, 1.1, fill=INDIGO)
txt(s6, "Feature 2 — AI Chat", 0.35, 0.12, 9, 0.5, size=22, bold=True, color=WHITE)
txt(s6, "Ask questions in plain English — get tables, insights, and analysis powered by GPT-4o", 0.35, 0.62, 10, 0.35, size=11, color=GRAY_300)
badge(s6, "GPT-4o POWERED", 10.3, 0.25, 2.3, 0.45, bg=INDIGO)

# Mock chat UI
rect(s6, 0.4, 1.2, 7.5, 5.9, fill=RGBColor(0x06,0x0D,0x20))
txt(s6, "AI Chat  |  FinanceData_MultiTab.xlsx", 0.55, 1.28, 6.0, 0.28, size=8, bold=True, color=WHITE)
badge(s6, "Connected", 6.6, 1.25, 1.1, 0.3, bg=RGBColor(0x05,0x46,0x2E), fg=EMERALD, size=7)

# Suggested questions strip
rect(s6, 0.55, 1.65, 7.2, 0.9, fill=SLATE_800)
txt(s6, "SUGGESTED QUESTIONS", 0.65, 1.68, 4.0, 0.22, size=7, bold=True, color=GRAY_500)
for i, q in enumerate(["Top 10 by Principal?", "Default customers table", "Credit score dist.", "Group by loan type"]):
    rect(s6, 0.65+i*1.75, 1.92, 1.65, 0.52, fill=RGBColor(0x0F,0x17,0x2A))
    txt(s6, q, 0.72+i*1.75, 1.98, 1.5, 0.38, size=7.5, color=GRAY_300)

# User message
rect(s6, 3.2, 2.72, 4.5, 0.48, fill=BLUE)
txt(s6, "Show top 10 loans by outstanding_balance", 3.35, 2.78, 4.2, 0.36, size=8.5, color=WHITE)

# AI table response
rect(s6, 0.55, 3.32, 7.2, 2.6, fill=SLATE_800)
txt(s6, "AI", 0.65, 3.38, 0.45, 0.38, size=9, bold=True,
    color=WHITE)
rect(s6, 0.65, 3.38, 0.45, 0.38, fill=BLUE)
txt(s6, "Top 10 Loans by Outstanding Balance", 1.2, 3.36, 6.0, 0.28, size=9, bold=True, color=WHITE)

# Table header
rect(s6, 0.65, 3.72, 7.0, 0.3, fill=RGBColor(0x1E,0x29,0x3B))
for j, h in enumerate(["LOAN_ID","CUSTOMER_ID","TYPE","AMOUNT","OUTSTANDING","STATUS"]):
    txt(s6, h, 0.72+j*1.12, 3.74, 1.05, 0.24, size=6.5, bold=True, color=BLUE)

# Table rows
row_data = [
    ("L00042","C00182","Mortgage","$500K","$478K","Active"),
    ("L00091","C00034","Business","$250K","$241K","Active"),
    ("L00017","C00156","Mortgage","$250K","$237K","Watch"),
    ("L00063","C00099","Personal","$100K","$98K","Default"),
    ("L00128","C00071","Auto","$50K","$49K","Active"),
]
for ri, row in enumerate(row_data):
    bg = RGBColor(0x0F,0x17,0x2A) if ri%2==0 else SLATE_800
    rect(s6, 0.65, 4.05+ri*0.36, 7.0, 0.34, fill=bg)
    for j, cell in enumerate(row):
        c = RED if cell == "Default" else AMBER if cell == "Watch" else WHITE
        txt(s6, cell, 0.72+j*1.12, 4.08+ri*0.36, 1.05, 0.26, size=7, color=c)

txt(s6, "Total: 10 loans shown | Combined outstanding: $1.23M | 2 require immediate attention (Watch/Default)",
    0.65, 5.95, 7.0, 0.35, size=7.5, bold=True, color=AMBER)

# Right: capabilities
caps = [
    ("Table Queries", "Show top/bottom N, filter by condition, group-by, cross-tab"),
    ("Smart Detection", "Detects question type — analytics vs tabular — automatically"),
    ("Real Data Rows", "Pre-filters actual DataFrame before sending to GPT-4o"),
    ("10 Smart Questions", "Dataset-specific questions generated on upload per tab"),
    ("GFM Tables", "Markdown tables rendered with zebra rows and hover highlights"),
    ("NLP Summary", "Every table answer includes a plain-English insight paragraph"),
]
txt(s6, "CAPABILITIES", 8.2, 1.22, 4.8, 0.3, size=8.5, bold=True, color=INDIGO)
for i, (t, d) in enumerate(caps):
    rect(s6, 8.2, 1.6+i*0.82, 4.8, 0.72, fill=SLATE_800)
    rect(s6, 8.2, 1.6+i*0.82, 0.04, 0.72, fill=INDIGO if i%2==0 else BLUE)
    txt(s6, t, 8.32, 1.62+i*0.82, 4.6, 0.3, size=9.5, bold=True, color=WHITE)
    txt(s6, d, 8.32, 1.9+i*0.82,  4.6, 0.3, size=8,   color=GRAY_300)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 7 — Feature: Analyst Workbench
# ═══════════════════════════════════════════════════════════════════════════════
s7 = slide()
rect(s7, 0, 0, 13.33, 1.1, fill=RGBColor(0x06,0x0B,0x18))
rect(s7, 0, 0, 0.06, 1.1, fill=EMERALD)
txt(s7, "Feature 3 — Analyst Workbench", 0.35, 0.12, 9, 0.5, size=22, bold=True, color=WHITE)
txt(s7, "Build any custom chart with point-and-click column selection — no code required", 0.35, 0.62, 10, 0.35, size=11, color=GRAY_300)
badge(s7, "SELF-SERVICE", 10.5, 0.25, 2.1, 0.45, bg=EMERALD)

# Workbench mock
rect(s7, 0.4, 1.2, 7.5, 5.9, fill=RGBColor(0x06,0x0D,0x20))
txt(s7, "Analyst Workbench  |  Loans tab", 0.55, 1.28, 6.0, 0.28, size=8, bold=True, color=WHITE)

rect(s7, 0.55, 1.65, 7.2, 0.7, fill=SLATE_800)
txt(s7, "CHART TYPE", 0.72, 1.72, 1.5, 0.22, size=7, bold=True, color=GRAY_500)
txt(s7, "X AXIS", 2.5, 1.72, 1.5, 0.22, size=7, bold=True, color=GRAY_500)
txt(s7, "Y AXIS", 4.3, 1.72, 1.5, 0.22, size=7, bold=True, color=GRAY_500)
rect(s7, 0.72, 1.96, 1.5, 0.3, fill=RGBColor(0x0F,0x17,0x2A))
rect(s7, 2.5, 1.96, 1.5, 0.3, fill=RGBColor(0x0F,0x17,0x2A))
rect(s7, 4.3, 1.96, 1.5, 0.3, fill=RGBColor(0x0F,0x17,0x2A))
txt(s7, "Bar Chart", 0.82, 1.98, 1.3, 0.24, size=8, color=WHITE)
txt(s7, "loan_type", 2.6, 1.98, 1.3, 0.24, size=8, color=WHITE)
txt(s7, "outstanding_balance", 4.4, 1.98, 1.5, 0.24, size=8, color=WHITE)
rect(s7, 6.1, 1.88, 1.5, 0.48, fill=BLUE)
txt(s7, "Generate Chart", 6.15, 1.96, 1.4, 0.32, size=8, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

rect(s7, 0.55, 2.48, 7.2, 4.3, fill=SLATE_800)
txt(s7, "Outstanding Balance by Loan Type", 0.72, 2.56, 5.0, 0.28, size=9, bold=True, color=WHITE)
bar_data = [("Personal",78),("Mortgage",92),("Auto",45),("Business",88),("Education",62),("Revolving",35)]
for i,(lbl,val) in enumerate(bar_data):
    bx = 0.72 + i*1.13
    cols = [BLUE,INDIGO,EMERALD,AMBER,RED,RGBColor(0x06,0xB6,0xD4)]
    rect(s7, bx, 6.3 - val/100*2.8, 0.9, val/100*2.8, fill=cols[i])
    txt(s7, lbl, bx, 6.35, 0.9, 0.25, size=6.5, color=GRAY_300, align=PP_ALIGN.CENTER)

# Right: chart types grid
txt(s7, "AVAILABLE CHART TYPES", 8.2, 1.22, 4.8, 0.3, size=8.5, bold=True, color=EMERALD)
charts = [("Bar","Compare categories"),("Scatter","Find correlations"),("Histogram","View distributions"),
          ("Box Plot","Detect outliers"),("Heatmap","Correlation matrix"),("Pie / Donut","Show proportions")]
for i, (ct, cd) in enumerate(charts):
    bx = 8.2 + (i%2)*2.5
    by = 1.6 + (i//2)*1.0
    rect(s7, bx, by, 2.3, 0.85, fill=SLATE_800)
    rect(s7, bx, by, 0.04, 0.85, fill=EMERALD if i%2==0 else BLUE)
    txt(s7, ct, bx+0.12, by+0.08, 2.1, 0.3, size=11, bold=True, color=WHITE)
    txt(s7, cd, bx+0.12, by+0.44, 2.1, 0.32, size=8.5, color=GRAY_300)

rect(s7, 8.2, 4.72, 4.8, 1.4, fill=RGBColor(0x05,0x23,0x1A))
rect(s7, 8.2, 4.72, 0.04, 1.4, fill=EMERALD)
txt(s7, "Self-Service Analytics", 8.32, 4.78, 4.6, 0.3, size=10, bold=True, color=EMERALD)
for ln in ["Any analyst can build production-quality charts",
           "without waiting for a data science team.",
           "Point, click, and generate — instantly."]:
    pass
txt(s7, "Any analyst can build production-quality charts\nwithout waiting for a data science team.\nPoint, click, and generate — instantly.",
    8.32, 5.1, 4.6, 0.9, size=9, color=GRAY_300, italic=True)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 8 — Business Use Cases
# ═══════════════════════════════════════════════════════════════════════════════
s8 = slide()
rect(s8, 0, 0, 13.33, 1.1, fill=RGBColor(0x06,0x0B,0x18))
rect(s8, 0, 0, 0.06, 1.1, fill=AMBER)
txt(s8, "Business Use Cases", 0.35, 0.12, 9, 0.5, size=24, bold=True, color=WHITE)
txt(s8, "Real-world scenarios where PulseAI delivers immediate value", 0.35, 0.62, 10, 0.35, size=11, color=GRAY_300)
badge(s8, "USE CASES", 10.7, 0.25, 1.9, 0.45, bg=AMBER, fg=RGBColor(0,0,0))

cases = [
    (BLUE,   "Banking & Financial Services", "Risk Management",
     "Upload Banking_LoanPortfolio.xlsx (5 tabs, 1,820 rows). Instantly identify\ndefault-risk customers, outstanding balance concentration, and DPD trends.",
     ["34% Default+Watch Rate detected", "Portfolio concentration risk", "Credit score distribution", "Cross-tab: loan type vs status"],
     [("34%","Default+Watch",RED),("$743.9M","Outstanding",BLUE),("657","Avg Credit",EMERALD)]),
    (EMERALD,"Insurance", "Claims Analytics",
     "Analyze Insurance_Claims.xlsx (4 tabs, 1,230 rows). Spot high-frequency\nclaimants, fraud indicators, and seasonal patterns without any SQL.",
     ["40% claims flagged for fraud", "180 fraud indicators found", "17.3% rejection rate", "61-day avg processing time"],
     [("40%","Fraud Flagged",RED),("180","Indicators",AMBER),("61 days","Avg Processing",EMERALD)]),
    (INDIGO, "Retail & E-Commerce", "Sales Performance",
     "Drop in Retail_Sales.xlsx (5 tabs, 3,062 rows). Get revenue KPIs,\ntop/bottom performers, regional breakdowns, and AI-written summaries.",
     ["21.4% return rate detected", "Top 3 categories nearly equal share", "$401.1M total line revenue", "16.4% cancelled order rate"],
     [("21.4%","Return Rate",RED),("$401.1M","Revenue",BLUE),("16.4%","Cancelled",AMBER)]),
    (AMBER,  "HR & Workforce", "People Analytics",
     "Upload HR_Workforce.xlsx (5 tabs, 3,411 rows). Identify attrition risks,\ncompensation outliers, and diversity metrics in minutes.",
     ["66.8% active employee rate", "5 departments at attrition risk", "1.09x gender pay gap detected", "$548.3M total payroll"],
     [("66.8%","Active Rate",EMERALD),("$548.3M","Payroll",BLUE),("1.09x","Pay Gap",AMBER)]),
]

for i, (col, domain, title, desc, bullets, metrics) in enumerate(cases):
    row = i // 2
    col_pos = i % 2
    bx = 0.4 + col_pos * 6.4
    by = 1.25 + row * 2.95
    rect(s8, bx, by, 6.1, 2.75, fill=SLATE_800)
    rect(s8, bx, by, 6.1, 0.05, fill=col)
    badge(s8, domain, bx+0.1, by+0.12, 2.5, 0.38, bg=RGBColor(0x06,0x0B,0x18), fg=col, size=8)
    txt(s8, title, bx+0.15, by+0.6, 5.8, 0.35, size=13, bold=True, color=WHITE)
    txt(s8, desc, bx+0.15, by+0.98, 3.5, 0.65, size=8, color=GRAY_300)
    # Metric tiles (3 small tiles below desc)
    for k, (mv, ml, mc) in enumerate(metrics):
        mx = bx+0.15 + k*1.95
        rect(s8, mx, by+1.72, 1.8, 0.72, fill=RGBColor(0x06,0x0B,0x18))
        txt(s8, mv, mx+0.08, by+1.76, 1.65, 0.34, size=13, bold=True, color=mc)
        txt(s8, ml, mx+0.08, by+2.08, 1.65, 0.28, size=7, color=GRAY_300)
    # Sample questions on right side
    for j, b in enumerate(bullets):
        txt(s8, f"-> {b}", bx+3.8, by+0.62+j*0.47, 2.2, 0.42, size=7.5, color=col)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 9 — ROI & Benefits
# ═══════════════════════════════════════════════════════════════════════════════
s9 = slide()
rect(s9, 0, 0, 13.33, 1.1, fill=RGBColor(0x06,0x0B,0x18))
rect(s9, 0, 0, 0.06, 1.1, fill=EMERALD)
txt(s9, "ROI & Business Benefits", 0.35, 0.12, 9, 0.5, size=24, bold=True, color=WHITE)
txt(s9, "Quantified value delivered to teams and leadership", 0.35, 0.62, 10, 0.35, size=11, color=GRAY_300)
badge(s9, "ROI", 11.5, 0.25, 1.1, 0.45, bg=EMERALD)

# Big stat row
for i, (num, unit, label, col) in enumerate([
    ("95%",  "faster",   "Time to insight vs manual Excel analysis", BLUE),
    ("Zero", "code",     "No SQL, Python, or BI tool training needed", EMERALD),
    ("< 2",  "minutes",  "From file upload to first actionable insight", AMBER),
    ("10x",  "more data","Questions answered per analyst per day", INDIGO),
]):
    bx = 0.4 + i * 3.2
    rect(s9, bx, 1.2, 3.0, 1.5, fill=SLATE_800)
    rect(s9, bx, 1.2, 3.0, 0.05, fill=col)
    txt(s9, num,   bx+0.15, 1.32, 2.7, 0.55, size=30, bold=True, color=col)
    txt(s9, unit,  bx+0.15, 1.85, 2.7, 0.3,  size=13, color=WHITE)
    txt(s9, label, bx+0.15, 2.2,  2.7, 0.38, size=8.5, color=GRAY_300)

# Benefits grid
benefits = [
    (BLUE,   "For Business Leaders",   ["Board-ready dashboards in minutes", "AI-written narrative — no analyst needed", "Confident decisions from real data", "All tabs in one unified view"]),
    (INDIGO, "For Data Analysts",      ["Skip manual chart building", "Self-service — no IT dependency", "Custom charts with 2 clicks", "Export-ready Plotly visualizations"]),
    (EMERALD,"For IT & Data Teams",    ["Zero infrastructure change", "Works with existing Excel/CSV files", "SQLite local DB — no data leaves your env", "Pluggable LLM — swap GPT-4o anytime"]),
    (AMBER,  "For Finance & Risk",     ["Instant loan portfolio risk view", "Default & outlier auto-detection", "Cross-tab drill-down without SQL", "Regulatory-ready audit tables"]),
]
for i, (col, title, pts) in enumerate(benefits):
    bx = 0.4 + (i%2)*6.4
    by = 3.1 + (i//2)*2.05
    rect(s9, bx, by, 6.1, 1.85, fill=SLATE_800)
    rect(s9, bx, by, 0.04, 1.85, fill=col)
    txt(s9, title, bx+0.15, by+0.1, 5.8, 0.32, size=11, bold=True, color=col)
    for j, pt in enumerate(pts):
        txt(s9, f"+ {pt}", bx+0.15, by+0.5+j*0.3, 5.8, 0.28, size=8.5, color=WHITE if j==0 else GRAY_300)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 10 — Technology Stack
# ═══════════════════════════════════════════════════════════════════════════════
s10 = slide()
rect(s10, 0, 0, 13.33, 1.1, fill=RGBColor(0x06,0x0B,0x18))
rect(s10, 0, 0, 0.06, 1.1, fill=INDIGO)
txt(s10, "Technology Stack", 0.35, 0.12, 9, 0.5, size=24, bold=True, color=WHITE)
txt(s10, "Enterprise-grade, cloud-native, production-ready components", 0.35, 0.62, 10, 0.35, size=11, color=GRAY_300)
badge(s10, "TECH STACK", 10.5, 0.25, 2.1, 0.45, bg=INDIGO)

categories = [
    ("Frontend", BLUE, [
        ("React 18", "UI framework with hooks & Zustand state"),
        ("TypeScript", "Type-safe component development"),
        ("Vite 5", "Lightning-fast dev server & bundler"),
        ("Tailwind CSS", "Utility-first dark-mode design system"),
        ("Plotly.js", "Interactive charts & visualizations"),
        ("react-markdown + GFM", "Formatted AI responses with tables"),
    ]),
    ("Backend", INDIGO, [
        ("FastAPI", "Async Python REST + WebSocket API"),
        ("SQLAlchemy 2.x", "ORM with SQLite (local) / PostgreSQL (prod)"),
        ("Uvicorn", "ASGI server — production-grade"),
        ("Pydantic v2", "Data validation & settings management"),
        ("pandas 2.2", "DataFrame processing & analysis"),
        ("openpyxl 3.1.5", "Excel file parsing (multi-tab support)"),
    ]),
    ("AI & Integration", EMERALD, [
        ("GPT-4o", "Frontier LLM for analysis & narrative"),
        ("Accenture lbpass", "Enterprise GenAI gateway (secure)"),
        ("WebSocket", "Real-time streaming chat connection"),
        ("remark-gfm", "GitHub-flavored markdown table rendering"),
        ("Custom LLM Factory", "Swappable provider abstraction"),
        ("Prompt Engineering", "Context-aware table & analytics prompts"),
    ]),
]

for i, (cat, col, items) in enumerate(categories):
    bx = 0.4 + i * 4.28
    rect(s10, bx, 1.2, 4.0, 5.9, fill=SLATE_800)
    rect(s10, bx, 1.2, 4.0, 0.05, fill=col)
    txt(s10, cat.upper(), bx+0.15, 1.3, 3.7, 0.32, size=10, bold=True, color=col)
    for j, (tech, desc) in enumerate(items):
        rect(s10, bx+0.12, 1.78+j*0.82, 3.76, 0.7, fill=RGBColor(0x0F,0x17,0x2A))
        txt(s10, tech, bx+0.25, 1.82+j*0.82, 3.5, 0.3, size=10, bold=True, color=WHITE)
        txt(s10, desc, bx+0.25, 2.1+j*0.82,  3.5, 0.28, size=8, color=GRAY_300)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 11 — Live Demo: Finance Dataset Example
# ═══════════════════════════════════════════════════════════════════════════════
s11 = slide()
rect(s11, 0, 0, 13.33, 1.1, fill=RGBColor(0x06,0x0B,0x18))
rect(s11, 0, 0, 0.06, 1.1, fill=AMBER)
txt(s11, "Live Example — Banking Loan Portfolio", 0.35, 0.12, 9, 0.5, size=22, bold=True, color=WHITE)
txt(s11, "Banking_LoanPortfolio.xlsx — 5 tabs, 1,820 rows, real FK relationships auto-detected", 0.35, 0.62, 10, 0.35, size=11, color=GRAY_300)
badge(s11, "BANKING DEMO", 10.7, 0.25, 1.9, 0.45, bg=AMBER, fg=RGBColor(0,0,0))

# Tab list
tabs = [
    ("Customers",    "200 rows",  "customer_id, full_name, age, segment, city, annual_income, kyc_status"),
    ("Loans",        "300 rows",  "loan_id, customer_id, loan_type, outstanding_balance, dpd, loan_status, credit_score"),
    ("Loan_Payments","900 rows",  "payment_id, loan_id, scheduled_emi, amount_paid, delay_days, payment_status"),
    ("Credit_Bureau","200 rows",  "customer_id, bureau_score, active_loans_count, credit_utilization_pct, dpd_30_count"),
    ("Collateral",   "220 rows",  "collateral_id, loan_id, collateral_type, market_value, ltv_ratio, status"),
]

txt(s11, "DATASET STRUCTURE", 0.5, 1.2, 5.0, 0.28, size=8.5, bold=True, color=AMBER)
for i, (tab, rows, cols) in enumerate(tabs):
    rect(s11, 0.5, 1.55+i*0.87, 6.0, 0.77, fill=SLATE_800)
    rect(s11, 0.5, 1.55+i*0.87, 0.04, 0.77, fill=AMBER if i==0 else BLUE)
    badge(s11, rows, 0.7, 1.6+i*0.87, 0.95, 0.28, bg=RGBColor(0x0F,0x17,0x2A), fg=AMBER, size=7)
    txt(s11, tab, 1.75, 1.58+i*0.87, 4.6, 0.3, size=10, bold=True, color=WHITE)
    txt(s11, cols, 1.75, 1.86+i*0.87, 4.6, 0.3, size=7.5, color=GRAY_300)

# Relationships
txt(s11, "FK RELATIONSHIPS DETECTED", 6.9, 1.2, 5.9, 0.28, size=8.5, bold=True, color=AMBER)
rels = [
    ("Loans",        "customer_id", "Customers",    "customer_id", "One customer → many loans"),
    ("Loan_Payments","loan_id",     "Loans",         "loan_id",    "One loan → many EMI payments"),
    ("Credit_Bureau","customer_id", "Customers",    "customer_id", "Credit profile per customer"),
    ("Collateral",   "loan_id",     "Loans",         "loan_id",    "Collateral pledged per loan"),
]
for i, (t1, c1, t2, c2, desc) in enumerate(rels):
    rect(s11, 6.9, 1.55+i*1.12, 6.0, 0.88, fill=SLATE_800)
    rect(s11, 6.9, 1.55+i*1.12, 0.04, 0.88, fill=INDIGO)
    txt(s11, f"{t1}.{c1}", 7.05, 1.6+i*1.12, 2.4, 0.28, size=9, bold=True, color=BLUE)
    txt(s11, "->", 9.5, 1.6+i*1.12, 0.4, 0.28, size=11, bold=True, color=AMBER)
    txt(s11, f"{t2}.{c2}", 9.95, 1.6+i*1.12, 2.4, 0.28, size=9, bold=True, color=INDIGO)
    txt(s11, desc, 7.05, 1.9+i*1.12, 5.7, 0.28, size=8, color=GRAY_300, italic=True)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 12 — Roadmap
# ═══════════════════════════════════════════════════════════════════════════════
s12 = slide()
rect(s12, 0, 0, 13.33, 1.1, fill=RGBColor(0x06,0x0B,0x18))
rect(s12, 0, 0, 0.06, 1.1, fill=BLUE)
txt(s12, "Roadmap & Next Steps", 0.35, 0.12, 9, 0.5, size=24, bold=True, color=WHITE)
txt(s12, "What's live today and where PulseAI goes next", 0.35, 0.62, 10, 0.35, size=11, color=GRAY_300)
badge(s12, "ROADMAP", 10.7, 0.25, 1.9, 0.45, bg=BLUE)

phases = [
    ("v1.0  NOW", EMERALD, [
        ("Done", "Multi-tab Excel + CSV upload & parsing"),
        ("Done", "Executive Dashboard with 6 auto-charts"),
        ("Done", "AI Chat with table & analytics answers"),
        ("Done", "Analyst Workbench self-service charts"),
        ("Done", "GPT-4o narrative + insight engine"),
        ("Done", "FK relationship auto-detection"),
    ]),
    ("v1.5  Q3 2026", BLUE, [
        ("Planned", "PDF & Google Sheets connector"),
        ("Planned", "Export charts to PowerPoint / PDF"),
        ("Planned", "Scheduled email reports"),
        ("Planned", "Multi-user workspaces & roles"),
        ("Planned", "Natural language chart builder"),
        ("Planned", "Data quality scoring dashboard"),
    ]),
    ("v2.0  Q4 2026", INDIGO, [
        ("Vision", "PostgreSQL & Snowflake connectors"),
        ("Vision", "Agent-based multi-step data queries"),
        ("Vision", "Real-time streaming data support"),
        ("Vision", "Fine-tuned domain LLM (Finance/HR)"),
        ("Vision", "Azure Container Apps deployment"),
        ("Vision", "SOC 2 & GDPR compliance module"),
    ]),
]
for i, (phase, col, items) in enumerate(phases):
    bx = 0.4 + i * 4.28
    rect(s12, bx, 1.2, 4.0, 5.7, fill=SLATE_800)
    rect(s12, bx, 1.2, 4.0, 0.06, fill=col)
    txt(s12, phase, bx+0.15, 1.3, 3.7, 0.35, size=12, bold=True, color=col)
    for j, (status, feat) in enumerate(items):
        sc = EMERALD if status=="Done" else BLUE if status=="Planned" else INDIGO
        rect(s12, bx+0.12, 1.78+j*0.77, 3.76, 0.65, fill=RGBColor(0x0F,0x17,0x2A))
        badge(s12, status, bx+0.18, 1.82+j*0.77, 0.75, 0.22, bg=RGBColor(0x06,0x0B,0x18), fg=sc, size=6.5)
        txt(s12, feat, bx+1.05, 1.82+j*0.77, 2.7, 0.52, size=8.5, color=WHITE)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 13 — Call to Action / Closing
# ═══════════════════════════════════════════════════════════════════════════════
s13 = slide()
rect(s13, 0, 0, 13.33, 7.5, fill=RGBColor(0x06,0x08,0x14))
rect(s13, 0, 0, 0.06, 7.5, fill=BLUE)

txt(s13, "Ready to transform\nyour data into decisions?",
    0.55, 0.7, 9.0, 1.5, size=30, bold=True, color=WHITE)

txt(s13, "PulseAI is live, running, and ready for your data today.",
    0.55, 2.35, 9.0, 0.45, size=14, color=BLUE)

txt(s13, "Upload any Excel or CSV file — and walk away with\nKPIs, charts, AI narrative, and actionable insights in under 2 minutes.",
    0.55, 2.9, 9.5, 0.75, size=12, color=GRAY_300)

# CTA buttons
rect(s13, 0.55, 3.9, 2.8, 0.65, fill=BLUE)
txt(s13, "Try Live Demo", 0.65, 3.97, 2.6, 0.5, size=12, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

rect(s13, 3.65, 3.9, 2.8, 0.65, fill=SLATE_800)
txt(s13, "View Architecture", 3.75, 3.97, 2.6, 0.5, size=12, bold=True, color=GRAY_300, align=PP_ALIGN.CENTER)

# Stack badges
for i, tech in enumerate(["React 18","FastAPI","GPT-4o","SQLite","Plotly","Tailwind"]):
    badge(s13, tech, 0.55+i*1.55, 5.0, 1.35, 0.38, bg=RGBColor(0x0F,0x17,0x2A), fg=GRAY_300, size=8)

divider(s13, 5.65)

txt(s13, "PulseAI  |  Enterprise Data Analytics Platform",
    0.55, 5.8, 7.0, 0.35, size=10, bold=True, color=WHITE)
txt(s13, "Built with React 18 + FastAPI + GPT-4o  |  Accenture GenAI Gateway",
    0.55, 6.18, 8.0, 0.3, size=9, color=GRAY_500)
txt(s13, "Confidential — For internal use and customer briefings only",
    0.55, 6.55, 8.0, 0.3, size=8.5, italic=True, color=GRAY_500)

# Right: summary stats
rect(s13, 10.2, 0.8, 2.6, 5.9, fill=SLATE_800)
rect(s13, 10.2, 0.8, 2.6, 0.05, fill=INDIGO)
txt(s13, "AT A GLANCE", 10.35, 0.92, 2.3, 0.28, size=8, bold=True, color=INDIGO)
for i, (val, lbl) in enumerate([
    ("5", "Product Features"),
    ("6", "Chart Types"),
    ("5", "Excel Tabs"),
    ("10", "AI Questions/Tab"),
    ("GPT-4o", "AI Engine"),
    ("< 2min", "Time to Insight"),
    ("Zero", "Code Required"),
    ("4", "Industry Use Cases"),
]):
    rect(s13, 10.3, 1.3+i*0.58, 2.3, 0.5, fill=RGBColor(0x0F,0x17,0x2A))
    txt(s13, val, 10.42, 1.33+i*0.58, 2.0, 0.26, size=11, bold=True, color=BLUE)
    txt(s13, lbl, 10.42, 1.57+i*0.58, 2.0, 0.22, size=7.5, color=GRAY_300)


# ── Save ───────────────────────────────────────────────────────────────────────
from pathlib import Path
out = Path(r"C:\Users\n.sureshmanikandan\Repo1\PulseAI\test_data\PulseAI_Presentation.pptx")
out.parent.mkdir(parents=True, exist_ok=True)
prs.save(str(out))
print(f"[OK] Saved: {out}")
print(f"     Slides: {len(prs.slides)}")
print("     Slide list:")
titles = [
    "01 - Hero / Title",
    "02 - Problem Statement",
    "03 - Solution Overview",
    "04 - System Architecture",
    "05 - Feature: Executive Dashboard",
    "06 - Feature: AI Chat",
    "07 - Feature: Analyst Workbench",
    "08 - Business Use Cases",
    "09 - ROI & Benefits",
    "10 - Technology Stack",
    "11 - Live Example (Finance Dataset)",
    "12 - Roadmap",
    "13 - Call to Action",
]
for t in titles:
    print(f"       {t}")
